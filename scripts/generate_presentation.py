#!/usr/bin/env python3
"""Generate technical presentation PowerPoint for the ApeRAG platform with rich cloud-native visual design."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Enhanced Color Palette (Blue Theme for ApeRAG) ──────────────────────────
BLUE = RGBColor(0x22, 0x63, 0xEB)
BLUE_LIGHT = RGBColor(0xBF, 0xDB, 0xFE)
BLUE_DARK = RGBColor(0x1E, 0x40, 0xAF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_LIGHT = RGBColor(0x2D, 0x37, 0x4A)
GRAY = RGBColor(0x4A, 0x55, 0x68)
GRAY_LIGHT = RGBColor(0x71, 0x7D, 0x8B)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFB)
LIGHT_BG_ACCENT = RGBColor(0xED, 0xF2, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GREEN_LIGHT = RGBColor(0x86, 0xEF, 0xAC)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
ORANGE_LIGHT = RGBColor(0xFB, 0xBF, 0x24)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LIGHT = RGBColor(0xDDD6, 0xFE)
RED = RGBColor(0xDC, 0x26, 0x26)
RED_LIGHT = RGBColor(0xFE, 0xCA, 0xCA)

# Cloud-native log colors
LOG_BG = RGBColor(0x1E, 0x1E, 0x1E)
LOG_TEXT = RGBColor(0xD4, 0xD4, 0xD4)
LOG_GREEN = RGBColor(0x98, 0xC3, 0x79)
LOG_YELLOW = RGBColor(0xED, 0xD4, 0x00)
LOG_BLUE = RGBColor(0x61, 0xAF, 0xEF)
LOG_CYAN = RGBColor(0x56, 0xB6, 0xC2)
LOG_RED = RGBColor(0xE0, 0x6C, 0x75)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_rich_decorative_bg(slide, primary_color=BLUE, accent_color=BLUE_LIGHT):
    """Add rich decorative background with gradient-like layers and cloud-native elements."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Large decorative corner circles
    corner_circles = [
        (0, 0, 3.0, BLUE_LIGHT, 0.75),
        (prs.slide_width - Inches(3.0), 0, 2.5, GREEN_LIGHT, 0.70),
        (0, prs.slide_height - Inches(2.5), 2.5, PURPLE_LIGHT, 0.75),
        (prs.slide_width - Inches(3.5), prs.slide_height - Inches(3.0), 3.0, BLUE_LIGHT, 0.70),
    ]
    
    for x, y, size, color, transparency in corner_circles:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(size), Inches(size))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        circle.fill.transparency = transparency
    
    # Sidebar accent strip
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height)
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = primary_color
    sidebar.line.fill.background()
    
    # Top decorative bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent_color
    top_bar.line.fill.background()
    
    # Bottom footer strip
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, prs.slide_height - Inches(0.25), prs.slide_width, Inches(0.25)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = DARK
    footer.line.fill.background()
    
    add_cloud_shapes(slide)
    add_decorative_lines(slide, primary_color)


def add_cloud_shapes(slide):
    """Add decorative cloud-native icon representations."""
    # Pod-like shapes
    pod_positions = [
        (prs.slide_width - Inches(1.5), Inches(0.5), 0.15, BLUE),
        (prs.slide_width - Inches(1.2), Inches(0.8), 0.12, PURPLE),
        (prs.slide_width - Inches(1.8), Inches(0.7), 0.13, GREEN),
    ]
    
    for x, y, size, color in pod_positions:
        pod = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(size), Inches(size))
        pod.fill.solid()
        pod.fill.fore_color.rgb = color
        pod.line.color.rgb = DARK
        pod.line.width = Pt(1.5)
        pod.fill.transparency = 0.25


def add_decorative_lines(slide, color):
    """Add decorative accent lines."""
    for y_pos in [Inches(1.0), Inches(6.5)]:
        line = slide.shapes.add_connector(
            1, Inches(0.5), y_pos, Inches(12.5), y_pos
        )
        line.line.color.rgb = color
        line.line.width = Pt(1)
        line.line.dash_style = 2


def add_cloud_native_logs(slide, left, top, width, height, log_entries=None):
    """Add a cloud-native log panel to a slide."""
    if log_entries is None:
        log_entries = [
            ("INFO", "2026-03-04T09:15:23Z", "aperag-api", "GraphRAG query: 'What is MCP?'", LOG_BLUE),
            ("INFO", "2026-03-04T09:15:24Z", "neo4j", "Traversed 3 hops, found 5 entities", LOG_GREEN),
            ("INFO", "2026-03-04T09:15:25Z", "llm-router", "Routed to Claude-3.5-Sonnet", LOG_CYAN),
        ]
    
    log_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    log_panel.fill.solid()
    log_panel.fill.fore_color.rgb = LOG_BG
    log_panel.line.color.rgb = BLUE
    log_panel.line.width = Pt(3)
    
    header = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.15))
    header_tf = header.text_frame
    header_tf.text = "📊 ApeRAG Enterprise Logs"
    header_tf.paragraphs[0].font.size = Pt(11)
    header_tf.paragraphs[0].font.color.rgb = LOG_CYAN
    header_tf.paragraphs[0].font.bold = True
    
    y_offset = top + Inches(0.3)
    for level, timestamp, service, message, color in log_entries:
        log_text = f"[{timestamp}] {level:8s} {service:15s} {message}"
        text_box = slide.shapes.add_textbox(
            left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.25)
        )
        text_frame = text_box.text_frame
        text_frame.text = log_text
        text_frame.paragraphs[0].font.name = "Courier New"
        text_frame.paragraphs[0].font.size = Pt(9)
        text_frame.paragraphs[0].font.color.rgb = color
        y_offset += Inches(0.3)


def title_bar(slide, title_text, subtitle_text=None, accent_color=BLUE):
    """Enhanced title bar with decorative elements."""
    bar_height = Inches(1.3)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(0.4), prs.slide_width - Inches(1.2), bar_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.color.rgb = accent_color
    shape.line.width = Pt(4)
    
    accent_stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(0.4), Inches(0.2), bar_height
    )
    accent_stripe.fill.solid()
    accent_stripe.fill.fore_color.rgb = accent_color
    accent_stripe.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(1.2)
    tf.margin_top = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = DARK
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(5)


def txt(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    """Styled text box."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def bullets(slide, left, top, width, height, items, size=16, color=DARK, bullet_color=BLUE):
    """Styled bullet points."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u25b8  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return tf


def box(slide, left, top, width, height, title, body, border=BLUE, fill_color=WHITE, accent=True):
    """Enhanced box with decorative accent."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border
    shape.line.width = Pt(3)
    
    if accent:
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(0.18)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = border
        accent_bar.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.3) if accent else Inches(0.2)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = border
    
    if body:
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY


def add_footer_branding(slide, text="ApeRAG Platform"):
    """Add footer branding strip."""
    footer_text = slide.shapes.add_textbox(
        Inches(0.6), prs.slide_height - Inches(0.2), 
        prs.slide_width - Inches(1.2), Inches(0.1)
    )
    tf = footer_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"\u2699\ufe0f  {text}"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rich_decorative_bg(s, BLUE, BLUE_LIGHT)

# Decorative circles behind title
for size, color, trans in [(5.0, BLUE_LIGHT, 0.6), (4.5, GREEN_LIGHT, 0.65), (4.0, PURPLE_LIGHT, 0.7)]:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.5), Inches(size), Inches(size))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.fill.transparency = trans
    c.line.fill.background()

txt(s, 1.5, 1.5, 10, 1.2, "ApeRAG", 56, BLUE, True)
txt(s, 1.5, 2.8, 10, 0.8, "Enterprise Graph RAG & Agentic Intelligence", 28, DARK)
txt(s, 1.5, 3.8, 10, 0.5, "Graph RAG  \u2022  Hybrid Retrieval  \u2022  MCP  \u2022  Semantic Cache  \u2022  Multi-LLM", 18, GRAY)
txt(s, 1.5, 5.0, 8, 0.5, "Jiale Lin", 24, DARK, True)
txt(s, 1.5, 5.5, 8, 0.4, "jeremykalilin@gmail.com  \u2022  linkedin.com/in/jiale-lin-ab03a4149", 15, GRAY)
txt(s, 1.5, 6.2, 8, 0.4, "Technical Presentation \u2014 25 min  \u2022  github.com/ljluestc/ApeRAG", 14, GRAY)
add_footer_branding(s, "ApeRAG Platform \u2022 Technical Presentation")

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — About Me
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rich_decorative_bg(s, BLUE, BLUE_LIGHT)
title_bar(s, "About Me", "Jiale Lin \u2014 Introduction", BLUE)
txt(s, 0.8, 1.8, 5.5, 0.5, "Background", 22, BLUE, True)
bullets(s, 0.8, 2.3, 5.5, 4.0, [
    "M.S. Computer Science \u2014 University of Colorado Boulder",
    "Staff SWE at Aviatrix: distributed backend, multi-cloud networking",
    "RAG platform architecture + grounding (Pinecone/FAISS/Chroma)",
    "Agentic orchestration: LangGraph, tool calling, audit trails",
    "Platform: Kubernetes, Terraform, ArgoCD, Prometheus/Grafana",
], 15, DARK, BLUE)
txt(s, 7.0, 1.8, 5.5, 0.5, "Contributions to ApeRAG", 22, BLUE, True)
bullets(s, 7.0, 2.3, 5.5, 4.0, [
    "Bug fixes: null-safety in query.py, ValueError formatting",
    "Typo fix: MEMOTY \u2192 MEMORY with backward-compat alias",
    "22 new unit tests (query models, connector adaptor, prompts)",
    "Technical presentation (PPTX + PDF redesign)",
    "Context engineering and RAG grounding analysis",
], 15, DARK, BLUE)
add_footer_branding(s)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — Problem Statement
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rich_decorative_bg(s, ORANGE, ORANGE_LIGHT)
title_bar(s, "The Problem", "Why an enterprise RAG platform?", ORANGE)
txt(s, 0.8, 1.8, 11, 0.8, "Enterprise knowledge is trapped in documents, graphs, and databases. LLMs hallucinate without grounded retrieval and relational understanding.", 17)
box(s, 0.8, 3.0, 3.6, 1.8, "Multi-Source Knowledge", "Docs, PDFs, code, graphs.\nNo single retrieval method\ncovers all structural shapes.", border=ORANGE)
box(s, 4.8, 3.0, 3.6, 1.8, "Hallucination Risk", "Pure LLMs fabricate answers.\nNeed multi-hop reasoning +\ncitation for enterprise trust.", border=RED)
box(s, 8.8, 3.0, 3.6, 1.8, "Production Gaps", "Demos lack scalability.\nNeed auth, rate limits, MCP,\nK8s deploy, and caching.", border=BLUE)
add_cloud_native_logs(s, Inches(0.8), Inches(5.3), Inches(11), Inches(1.5), [
    ("ERROR", "2026-03-04T09:10:00Z", "llm-generator", "Hallucinated response: 'API v5' (does not exist)", LOG_RED),
    ("WARN", "2026-03-04T09:12:00Z", "retrieval", "Missing relational context for query 'A -> C'", LOG_YELLOW),
])
add_footer_branding(s)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — Architecture
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rich_decorative_bg(s, BLUE, BLUE_LIGHT)
title_bar(s, "System Architecture", "End-to-end Enterprise RAG Pipeline", BLUE)
stages = [
    ("Ingestion", "MinerU parsing\nSmart chunking\nGraph extraction", BLUE),
    ("Indexing", "Vector (Qdrant)\nFull-text (ES)\nGraph (Neo4j)", GREEN),
    ("Retrieval", "Vector + BM25\nGraph + Summary\nRe-ranking", PURPLE),
    ("Agentic Ops", "MCP Tool Calling\nSemantic Cache\nContent Guard", ORANGE),
    ("LLM Stack", "LiteLLM Router\nOpenAI/Anthropic\nLocal (Qwen/Llama)", BLUE),
]
x = 0.5
for t, b, c in stages:
    box(s, x, 2.0, 2.3, 2.2, t, b, border=c)
    x += 2.5
txt(s, 0.8, 4.8, 11, 0.5, "Infrastucture Layer", 18, BLUE, True)
bullets(s, 0.8, 5.3, 11, 1.5, [
    "Celery + Redis: async ingestion & background training workers",
    "ZooKeeper: distributed coordination & configuration management",
    "OpenTelemetry + Jaeger: distributed tracing across all microservices",
], 14, DARK, BLUE)
add_footer_branding(s)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — Enterprise Components (New!)
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rich_decorative_bg(s, PURPLE, PURPLE_LIGHT)
title_bar(s, "Enterprise RAG Infrastructure", "System resilience & scalability", PURPLE)
box(s, 0.8, 2.0, 3.6, 2.2, "API Gateway", "Unified entry point with\nAuth, Rate Limiting,\nand Request Routing.", border=BLUE)
box(s, 4.8, 2.0, 3.6, 2.2, "Semantic Cache", "Redis-backed cache that\nmatches query intent,\nreducing LLM costs.", border=GREEN)
box(s, 8.8, 2.0, 3.6, 2.2, "Content Moderation", "Safety filters for both\nprompt injection and\nharmful output detection.", border=RED)
box(s, 0.8, 4.5, 3.6, 1.8, "ZooKeeper", "Cluster state management\nand service discovery for\ndistributed workers.", border=ORANGE)
box(s, 4.8, 4.5, 3.6, 1.8, "Pub-Sub Messaging", "Kafka/Redis for async\ndocument events and\nstatus notifications.", border=PURPLE)
box(s, 8.8, 4.5, 3.6, 1.8, "Eval Pipeline", "Automated Ragas scoring\nintegrated into the\ndeployment lifecycle.", border=BLUE)
add_footer_branding(s)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — Five Index Types
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rich_decorative_bg(s, BLUE, BLUE_LIGHT)
title_bar(s, "Five Index Types", "Multi-dimensional understanding", BLUE)
types = [
    ("Vector", "High semantic recall\nCosine similarity", BLUE),
    ("Full-text", "Exact term matching\nBM25 ranking", GREEN),
    ("Graph", "Relational reasoning\nEntity extraction", PURPLE),
    ("Summary", "Broad document view\nGlobal context", ORANGE),
    ("Vision", "Multimodal support\nChart/Table OCR", RED),
]
x = 0.5
for t, b, c in types:
    box(s, x, 2.0, 2.3, 2.5, t, b, border=c)
    x += 2.5
txt(s, 0.8, 5.0, 11, 1.0, 
    "Hybrid retrieval merges these signals in parallel. BM25 catches exact identifiers, "
    "while Vector catches meaning, and Graph resolves relationships.", 15, GRAY)
add_cloud_native_logs(s, Inches(0.8), Inches(6.0), Inches(11), Inches(1.0), [
    ("INFO", "2026-03-04T09:15:30Z", "hybrid-search", "Merging Vector (0.8) + BM25 (0.5) + Graph (found path)", LOG_GREEN),
])
add_footer_branding(s)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — Enhanced Graph RAG
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rich_decorative_bg(s, GREEN, GREEN_LIGHT)
title_bar(s, "Enhanced Graph RAG", "LightRAG with entity normalization", GREEN)
txt(s, 0.8, 1.8, 5.5, 0.5, "KG Construction", 20, GREEN, True)
bullets(s, 0.8, 2.3, 5.5, 3.0, [
    "LLM-based entity & relation extraction",
    "Entity normalization to prevent duplicates",
    "Neo4j/NebulaGraph storage backend",
    "Graph visualization in Admin dashboard",
], 15, DARK, GREEN)
txt(s, 7.0, 1.8, 5.5, 0.5, "The Graph Advantage", 20, GREEN, True)
bullets(s, 7.0, 2.3, 5.5, 3.0, [
    "Multi-hop reasoning (A \u2192 B \u2192 C)",
    "Global context understanding",
    "Cross-document relationship mapping",
    "Entity-centric retrieval augmentation",
], 15, DARK, GREEN)
add_footer_branding(s)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 — AI Agents & MCP
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rich_decorative_bg(s, BLUE, BLUE_LIGHT)
title_bar(s, "AI Agents & MCP Integration", "Autonomous search and tool use", BLUE)
box(s, 0.8, 2.0, 3.6, 2.5, "Agent Orchestration", "Lifecycle management via\nAgentSessionManager.\nReal-time event streaming\nfor thinking process.", border=BLUE)
box(s, 4.8, 2.0, 3.6, 2.5, "MCP Protocol", "Universal interface for\nconnecting LLMs to\nknowledge bases and\nexternal API tools.", border=GREEN)
box(s, 8.8, 2.0, 3.6, 2.5, "Search Fallbacks", "Integrated web search\n(DuckDuckGo) when internal\nknowledge is insufficient\nor out of date.", border=ORANGE)
add_cloud_native_logs(s, Inches(0.8), Inches(5.3), Inches(11), Inches(1.5), [
    ("INFO", "2026-03-04T09:20:00Z", "agent", "Plan: Search KB \u2192 Tool: calculate \u2192 Final Answer", LOG_BLUE),
    ("INFO", "2026-03-04T09:20:05Z", "mcp-server", "Exporting context to Claude Desktop", LOG_CYAN),
])
add_footer_branding(s)

# ═══════════════════════════════════════════════════════════
# SLIDE 9 — Security & Observability
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rich_decorative_bg(s, RED, RED_LIGHT)
title_bar(s, "Security & Observability", "Production trust and transparency", RED)
txt(s, 0.8, 1.8, 5.5, 0.5, "Security & Auth", 20, RED, True)
bullets(s, 0.8, 2.3, 5.5, 3.0, [
    "JWT & OAuth2 (Google/GitHub)",
    "Role-based access control (RBAC)",
    "Rate limiting & audit logging",
    "Pydantic v2 input validation",
], 15, DARK, RED)
txt(s, 7.0, 1.8, 5.5, 0.5, "Observability Stack", 20, BLUE, True)
bullets(s, 7.0, 2.3, 5.5, 3.0, [
    "OpenTelemetry Auto-instrumentation",
    "Opik & Ragas for RAG evaluation",
    "Jaeger distributed tracing",
    "Prometheus metrics + Grafana",
], 15, DARK, BLUE)
add_footer_branding(s)

# ═══════════════════════════════════════════════════════════
# SLIDE 10 — Deployment
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rich_decorative_bg(s, BLUE, BLUE_LIGHT)
title_bar(s, "Deployment & Tech Stack", "Cloud-native from day one", BLUE)
box(s, 0.8, 2.0, 3.6, 2.0, "Infrastructure", "Docker Compose / Helm\nKubernetes (HPA support)\nKubeBlocks Database Ops", border=BLUE)
box(s, 4.8, 2.0, 3.6, 2.0, "Backend Stack", "Python 3.11, FastAPI\nCelery, Redis, Postgres\nLangChain, LiteLLM", border=GREEN)
box(s, 8.8, 2.0, 3.6, 2.0, "Frontend Stack", "Next.js, TypeScript\nTailwind CSS, ShadcnUI\nGraph-recharts vis", border=PURPLE)
add_cloud_native_logs(s, Inches(0.8), Inches(4.5), Inches(11), Inches(2.0), [
    ("INFO", "2026-03-04T09:30:00Z", "k8s-operator", "Scanning deployments: All pods healthy", LOG_GREEN),
    ("INFO", "2026-03-04T09:30:05Z", "helm-deploy", "Version 2.5.0 rollout complete", LOG_BLUE),
    ("INFO", "2026-03-04T09:30:10Z", "prometheus", "Recording SLIs: Latency P99 < 300ms", LOG_YELLOW),
])
add_footer_branding(s)

# ═══════════════════════════════════════════════════════════
# SLIDE 11 — My Contributions
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rich_decorative_bg(s, BLUE, BLUE_LIGHT)
title_bar(s, "My Contributions", "Hardware-accelerated RAG quality", BLUE)
txt(s, 0.8, 1.8, 5.5, 0.5, "Bug Fixes", 20, BLUE, True)
bullets(s, 0.8, 2.3, 5.5, 2.5, [
    "Fixed query.py None-safety crash",
    "Resolved connector adaptor formatting",
    "Fixed legacy MEMOTY constants",
], 15, DARK, BLUE)
txt(s, 7.0, 1.8, 5.5, 0.5, "New Features & Tests", 20, BLUE, True)
bullets(s, 7.0, 2.3, 5.5, 2.5, [
    "22 new unit tests for query logic",
    "Enterprise infrastructure planning",
    "Production-grade technical docs",
], 15, DARK, BLUE)
add_footer_branding(s)

# ═══════════════════════════════════════════════════════════
# SLIDE 12 — Conclusion & QA
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rich_decorative_bg(s, BLUE, BLUE_LIGHT)
title_bar(s, "Key Takeaways & Q\u200a&\u200aA", "Future-proof RAG platform", BLUE)
txt(s, 0.8, 1.8, 11, 2.0, 
    "ApeRAG combines Graph reasoning with Vector recall in a scalable enterprise architecture.\n\n"
    "Thank you for your time!", 24, DARK, True, align=PP_ALIGN.CENTER)
bullets(s, 1.5, 4.0, 10, 2.0, [
    "5 Index Types for maximum coverage",
    "Enterprise-ready with API Gateway & Semantic Cache",
    "Native MCP + Agent support for autonomous workflows",
    "Full observability with OpenTelemetry & Evaluators",
], 18, DARK, BLUE)
add_footer_branding(s, "Jiale Lin \u2022 ApeRAG Technical Presentation \u2022 QA session")

# ═══════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════
out_dir = os.path.join(os.path.dirname(__file__), "..", "docs")
os.makedirs(out_dir, exist_ok=True)
pptx_path = os.path.join(out_dir, "ApeRAG_Technical_Presentation.pptx")
prs.save(pptx_path)
print(f"Saved PPTX: {pptx_path}")
